from __future__ import annotations

import json
import os
from dataclasses import dataclass, field
from pathlib import Path

from .llm import (
    LLMClient,
    LLMConfigurationError,
    LLMRequestError,
    OpenAICompatibleClient,
    StreamEvent,
)
from .memory import MemoryStore
from .settings import LLMSettings, load_effective_llm_config
from .teaching import (
    TeachingSession,
    classify_teaching_input,
    clear_session,
    deconstruct_messages,
    has_active_session,
    load_session,
    parse_tree,
    render_tree,
    save_session,
    teach_messages,
)
from .workspace import SUPPORTED_SOURCE_EXTENSIONS, CramWorkspace, discover_workspace_sources
from .workspace_ingest import MaterialIngestResult, ingest_material_sources
from .workspace_index import (
    ChunkRecord,
    index_text_sources,
    load_workspace_chunks,
    search_workspace_chunks,
    workspace_chunks_path,
)


ARTIFACT_COMMANDS = {
    "/plan": (
        "速成计划.md",
        "期末速成计划",
        "输出一份考前速成路线：按优先级列出必须先掌握的考点，给出建议的学习顺序与时间分配，突出老师强调和高频考点。",
    ),
    "/notes": (
        "知识点整合.md",
        "知识点整合",
        "把分散在各份资料里的同一知识点整合成结构化笔记，按主题归类，每个知识点给出定义、关键要点和易混淆点。",
    ),
    "/mindmap": (
        "思维导图.md",
        "思维导图",
        "用 Markdown 多级无序列表输出思维导图，体现 章节 → 概念 → 关系 → 易混点 → 高频考点 的层级结构，结构优先于细节。",
    ),
    "/quiz": (
        "题库.md",
        "题库",
        "按常见考试题型出题（如选择、判断、简答、计算），覆盖重点考点；每题给出参考答案和踩分点。",
    ),
    "/summary": (
        "考前总结.md",
        "考前总结",
        "输出考前冲刺总结：必背要点清单、最易出错的点、以及考前最后一天的复习清单。",
    ),
}

HELP_TEXT = """可用命令：
/ingest   解析并索引当前文件夹资料
/status   查看资料、记忆和输出状态
/plan     生成速成计划
/notes    生成知识点整合
/mindmap  生成思维导图
/quiz     生成题库
/summary  生成考前总结
/lint     检查记忆、输出和引用冲突
/config   重新配置 LLM
/model    切换对话模型
/help     查看命令

直接输入问题即可继续复习对话。
"""

MAX_AGENT_STEPS = 8
# How many prior conversation turns (user+assistant pairs) to replay into the context window.
HISTORY_TURNS = 10

# Maps a tool name the model can call to the slash-style artifact key it generates.
TOOL_ARTIFACTS = {
    "generate_plan": "/plan",
    "generate_notes": "/notes",
    "generate_mindmap": "/mindmap",
    "generate_quiz": "/quiz",
    "generate_summary": "/summary",
}

_TOOL_LABELS = {
    "generate_plan": "生成速成计划",
    "generate_notes": "生成知识点整合",
    "generate_mindmap": "生成思维导图",
    "generate_quiz": "生成题库",
    "generate_summary": "生成考前总结",
    "ingest_materials": "导入并索引资料",
    "show_status": "查看状态",
    "lint_conflicts": "检查冲突",
    "start_teaching": "开始系统教学",
    "list_files": "查看文件列表",
    "read_file": "读取文件",
    "grep_materials": "搜索资料",
    "switch_workspace": "切换课程文件夹",
    "update_memory": "记入长期记忆",
}


def _artifact_tool(name: str, description: str) -> dict:
    return {
        "type": "function",
        "function": {
            "name": name,
            "description": description,
            "parameters": {
                "type": "object",
                "properties": {
                    "focus": {
                        "type": "string",
                        "description": '可选，本次生成要聚焦的主题/章节/题型，例如 "傅里叶变换" 或 "只出计算题"。',
                    }
                },
            },
        },
    }


def _plain_tool(name: str, description: str) -> dict:
    return {
        "type": "function",
        "function": {
            "name": name,
            "description": description,
            "parameters": {"type": "object", "properties": {}},
        },
    }


AGENT_TOOLS = [
    _artifact_tool("generate_plan", "生成或更新考前速成计划，写入 cram-output/速成计划.md。"),
    _artifact_tool("generate_notes", "把资料整合成结构化知识点笔记，写入 cram-output/知识点整合.md。"),
    _artifact_tool("generate_mindmap", "生成思维导图（Markdown 层级），写入 cram-output/思维导图.md。"),
    _artifact_tool("generate_quiz", "出题/生成题库或练习题，写入 cram-output/题库.md。"),
    _artifact_tool("generate_summary", "生成考前冲刺总结/必背清单，写入 cram-output/考前总结.md。"),
    _plain_tool("ingest_materials", "扫描并用 MinerU 解析、索引当前课程文件夹里的资料（PDF/PPT/图片等）。用户要导入/解析/重新索引资料时调用。"),
    _plain_tool("show_status", "查看当前课程的资料数量、产物数量、记忆与索引状态。"),
    _plain_tool("lint_conflicts", "检查长期记忆、生成产物与原始资料之间的引用冲突。"),
    {
        "type": "function",
        "function": {
            "name": "start_teaching",
            "description": "当用户想被系统地讲解/带着复习某个主题（如「教我X」「带我复习X」「我想系统学X」）时调用，进入四步教学法的互动教学（拆解知识点→逐点讲授）。只是单点答疑不要调用。",
            "parameters": {
                "type": "object",
                "properties": {
                    "topic": {
                        "type": "string",
                        "description": '要系统教学的主题，例如 "期望理论" 或 "第3章 信道编码"。',
                    }
                },
                "required": ["topic"],
            },
        },
    },
    _plain_tool("list_files", "列出当前课程文件夹里的资料文件和已生成的产物（cram-output）。用户想看有哪些文件时调用。"),
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "读取课程文件夹内某个文件的正文。md/txt 直接读；pdf/ppt/图片读其 MinerU 解析出的 markdown（没解析会提示先 /ingest）。只能读课程文件夹内的文件。",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "相对课程文件夹的文件路径，例如 老师重点.pdf 或 cram-output/速成计划.md。"}
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "grep_materials",
            "description": "在所有文本资料（md/txt 与已解析的 markdown）里按关键词搜索，返回命中文件、行号和该行内容。需要在原文里定位某个词出现在哪时调用。",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "要搜索的关键词或短语。"}
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "update_memory",
            "description": (
                "把一条值得长期记住的事实写入课程长期记忆（跨会话保留）。"
                "适合：考试形式（开卷/闭卷/题型）、老师强调的必考点、用户偏好、"
                "用户反复搞错或已经掌握的知识点。只记durable的事实，不要记一次性闲聊。"
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "note": {"type": "string", "description": "要记住的一句话事实，尽量具体。"},
                    "category": {
                        "type": "string",
                        "description": "可选分类：考点 / 偏好 / 易错 / 掌握 / 其他。",
                    },
                },
                "required": ["note"],
            },
        },
    },
]


def _api_tool_call(call: dict) -> dict:
    return {
        "id": call.get("id") or call.get("name", ""),
        "type": "function",
        "function": {
            "name": call.get("name", ""),
            "arguments": call.get("arguments") or "{}",
        },
    }


@dataclass(frozen=True)
class CommandResult:
    kind: str
    message: str
    wrote: list[Path] = field(default_factory=list)


class CommandRouter:
    def __init__(self, workspace: CramWorkspace, *, llm: LLMClient | None = None, material_ingester=None):
        self.workspace = workspace
        self.memory = MemoryStore.open(workspace)
        self.llm = llm or _default_llm_client()
        self.material_ingester = material_ingester or ingest_material_sources

    def handle(self, text: str) -> CommandResult:
        message = text.strip()
        if not message:
            return CommandResult(kind="empty", message="")

        self.memory.append_session_event("user", message)

        command = message.split(maxsplit=1)[0].lower()
        if command == "/help":
            return self._remember(CommandResult(kind="help", message=HELP_TEXT))
        if command == "/status":
            return self._remember(self._status())
        if command == "/lint":
            return self._remember(self._lint())
        if command == "/config":
            return CommandResult(kind="config", message="")
        if command == "/model":
            return CommandResult(kind="model", message="")
        if command == "/ingest":
            return self._remember(self._ingest_status())
        if command in ARTIFACT_COMMANDS:
            return self._remember(self._write_artifact(command))

        return self._remember(self._ask_llm(message))

    def _remember(self, result: CommandResult) -> CommandResult:
        if result.message:
            self.memory.append_session_event("agent", result.message)
        return result

    def _status(self) -> CommandResult:
        sources = discover_workspace_sources(self.workspace.root)
        outputs = [path for path in self.workspace.output_dir.rglob("*") if path.is_file()]
        references = self.memory.build_reference_catalog()
        message = (
            f"{self.workspace.course_name}\n"
            f"- {len(sources)} 个资料文件\n"
            f"- {len(outputs)} 个输出产物\n"
            f"- {len(references)} 个可引用条目\n"
            f"- 记忆目录：{self.workspace.cram_dir.relative_to(self.workspace.root).as_posix()}\n"
            f"- 输出目录：{self.workspace.output_dir.relative_to(self.workspace.root).as_posix()}"
        )
        return CommandResult(kind="status", message=message)

    def _ingest_status(self) -> CommandResult:
        sources = discover_workspace_sources(self.workspace.root)
        material_result: MaterialIngestResult = self.material_ingester(self.workspace, sources)
        index_result = index_text_sources(self.workspace, extra_texts=material_result.parsed_texts)
        manifest = self.workspace.cram_dir / "raw_manifest.json"
        manifest.write_text(
            "\n".join(source.relative_path.as_posix() for source in sources),
            encoding="utf-8",
        )
        return CommandResult(
            kind="ingest",
            message=_format_ingest_message(
                workspace=self.workspace,
                total_sources=len(sources),
                material_result=material_result,
                indexed_chunks=index_result.indexed_chunks,
                indexed_files=index_result.indexed_files,
            ),
        )

    def _write_artifact(self, command: str, *, focus: str | None = None) -> CommandResult:
        filename, title, instruction = ARTIFACT_COMMANDS[command]
        output_path = self.workspace.output_dir / filename
        evidence = self._collect_artifact_evidence(focus=focus)
        try:
            body = self.llm.chat(
                self._artifact_messages(title, instruction, evidence, focus=focus), stream=False
            )
        except LLMConfigurationError as exc:
            return CommandResult(kind="artifact", message=_llm_setup_message(exc))
        except LLMRequestError as exc:
            return CommandResult(
                kind="artifact",
                message=(
                    "生成失败，未写入文件。请检查模型服务地址、模型名和网络状态。\n"
                    f"{exc}"
                ),
            )
        output_path.write_text(self._format_artifact_file(body, evidence), encoding="utf-8")
        return CommandResult(
            kind="artifact",
            message=f"Wrote {output_path.relative_to(self.workspace.root).as_posix()}",
            wrote=[output_path],
        )

    def _collect_artifact_evidence(self, *, focus: str | None = None, limit_chars: int = 12000) -> list[ChunkRecord]:
        if focus:
            targeted = search_workspace_chunks(self.workspace, focus, limit=8)
            if targeted:
                return targeted
        selected: list[ChunkRecord] = []
        total = 0
        for chunk in load_workspace_chunks(self.workspace):
            if selected and total + len(chunk.text) > limit_chars:
                break
            selected.append(chunk)
            total += len(chunk.text)
        return selected

    def _artifact_messages(
        self, title: str, instruction: str, evidence: list[ChunkRecord], *, focus: str | None = None
    ) -> list[dict]:
        system_prompt = f"""你是期末速成引擎，为课程「{self.workspace.course_name}」生成可复用的复习产物：{title}。

生成原则：
- 资料优先：结论、定义、公式、考点尽量来自下方课程资料。
- 来源优先：用到资料的地方，在句末用方括号标注来源标签，例如 [a.pdf:mineru:1]。
- 整合优先：跨文件合并同一知识点，不要按文件机械罗列。
- 找不到依据就说明：资料中没有的内容标注为「推理补充」，不要编造来源或页码。
- 直接输出 Markdown 正文，不要寒暄或自我介绍。

本次任务：{instruction}
"""
        if focus:
            system_prompt += f"\n用户对本次生成的具体要求（请重点聚焦）：{focus}\n"
        if evidence:
            evidence_block = "\n\n".join(
                f"[{chunk.citation_label}]\n{chunk.text}" for chunk in evidence
            )
            system_prompt += f"\n课程资料片段：\n\n{evidence_block}\n"
        else:
            system_prompt += (
                "\n当前没有已索引的课程资料。请基于通用知识生成，"
                "并在开头用一句话说明：未检索到课程资料，建议先运行 /ingest 导入资料。\n"
            )
        user_request = f"请生成《{title}》。" + (f"重点：{focus}。" if focus else "")
        return [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_request},
        ]

    @staticmethod
    def _format_artifact_file(body: str, evidence: list[ChunkRecord]) -> str:
        if evidence:
            labels = "\n".join(f"- {chunk.citation_label}" for chunk in evidence)
            provenance = f"## 参考资料\n\n{labels}\n"
        else:
            provenance = "## 参考资料\n\n- 未检索到课程资料（建议先运行 /ingest）\n"
        return f"{body.strip()}\n\n---\n\n{provenance}"

    def _ask_llm(self, message: str) -> CommandResult:
        try:
            answer = self.llm.chat(self._llm_messages(message), stream=False)
        except LLMConfigurationError as exc:
            answer = _llm_setup_message(exc)
        except LLMRequestError as exc:
            answer = (
                "LLM 请求失败，当前会话已保留。请检查模型服务地址、模型名和网络状态。\n"
                f"{exc}"
            )
        return CommandResult(kind="ask", message=answer)

    def stream(self, message: str):
        self.memory.append_session_event("user", message)
        if not hasattr(self.llm, "stream_chat"):
            result = self._remember(self._ask_llm(message))
            if result.message:
                yield StreamEvent("content", result.message)
            return

        content_parts: list[str] = []
        try:
            for event in self.llm.stream_chat(self._llm_messages(message)):
                if not isinstance(event, StreamEvent):
                    event = StreamEvent("content", str(event))
                if event.kind == "content":
                    content_parts.append(event.text)
                yield event
        except LLMConfigurationError as exc:
            message_text = _llm_setup_message(exc)
            self.memory.append_session_event("agent", message_text)
            yield StreamEvent("content", message_text)
        except LLMRequestError as exc:
            message_text = (
                "LLM 请求失败，当前会话已保留。请检查模型服务地址、模型名和网络状态。\n"
                f"{exc}"
            )
            self.memory.append_session_event("agent", message_text)
            yield StreamEvent("content", message_text)
        else:
            answer = "".join(content_parts)
            if answer:
                self.memory.append_session_event("agent", answer)

    def run_turn(self, message: str):
        """Agentic turn: the model may call tools (generate/ingest/status/lint) or answer.

        Yields StreamEvent of kind reasoning/content/tool/wrote for the TUI. Falls back to a
        plain streamed answer when the client cannot do tool calls.
        """
        if has_active_session(self.workspace):
            yield from self._teaching_turn(message)
            return
        if not hasattr(self.llm, "stream_agent"):
            yield from self.stream(message)
            return

        self.memory.append_session_event("user", message)
        messages = self._agent_messages(message)
        answer_parts: list[str] = []
        try:
            for _ in range(MAX_AGENT_STEPS):
                assistant_content: list[str] = []
                tool_calls: list[dict] = []
                for event in self.llm.stream_agent(messages, AGENT_TOOLS):
                    if event.kind == "reasoning":
                        yield event
                    elif event.kind == "content":
                        assistant_content.append(event.text)
                        answer_parts.append(event.text)
                        yield event
                    elif event.kind == "tool_call":
                        try:
                            tool_calls.append(json.loads(event.text))
                        except ValueError:
                            continue
                    elif event.kind == "usage":
                        yield event
                if not tool_calls:
                    break
                messages.append(
                    {
                        "role": "assistant",
                        "content": "".join(assistant_content),
                        "tool_calls": [_api_tool_call(call) for call in tool_calls],
                    }
                )
                direct_text: str | None = None
                for call in tool_calls:
                    name = call.get("name", "")
                    yield StreamEvent("tool", _TOOL_LABELS.get(name, name))
                    if name == "switch_workspace":
                        result_text, switched = self._apply_switch_workspace(call)
                        if switched is not None:
                            yield StreamEvent("switch", str(switched))
                        wrote: list[Path] = []
                        present_directly = False
                    else:
                        result_text, wrote, present_directly = self._execute_tool(call)
                    for path in wrote:
                        yield StreamEvent("wrote", str(path))
                    if present_directly:
                        direct_text = result_text
                    messages.append(
                        {
                            "role": "tool",
                            "tool_call_id": call.get("id") or name,
                            "content": result_text,
                        }
                    )
                if direct_text is not None:
                    answer_parts.append(direct_text)
                    yield StreamEvent("content", direct_text)
                    break
            else:
                # Hit the tool-call cap without a final answer; force a closing reply (no tools).
                if not "".join(answer_parts).strip() and hasattr(self.llm, "stream_chat"):
                    messages.append(
                        {"role": "user", "content": "请基于上面的结果，用一两句话直接回复我，不要再调用工具。"}
                    )
                    for event in self.llm.stream_chat(messages):
                        if not isinstance(event, StreamEvent):
                            event = StreamEvent("content", str(event))
                        if event.kind == "content":
                            answer_parts.append(event.text)
                        yield event
        except LLMConfigurationError as exc:
            text = _llm_setup_message(exc)
            self.memory.append_session_event("agent", text)
            yield StreamEvent("content", text)
            return
        except LLMRequestError as exc:
            text = (
                "LLM 请求失败，当前会话已保留。请检查模型服务地址、模型名和网络状态。\n"
                f"{exc}"
            )
            self.memory.append_session_event("agent", text)
            yield StreamEvent("content", text)
            return

        answer = "".join(answer_parts).strip()
        if answer:
            self.memory.append_session_event("agent", answer)

    def _execute_tool(self, call: dict) -> tuple[str, list[Path], bool]:
        name = call.get("name", "")
        try:
            args = json.loads(call.get("arguments") or "{}")
        except ValueError:
            args = {}
        if not isinstance(args, dict):
            args = {}

        if name in TOOL_ARTIFACTS:
            focus = args.get("focus") or None
            result = self._write_artifact(TOOL_ARTIFACTS[name], focus=focus)
            if result.wrote:
                title = ARTIFACT_COMMANDS[TOOL_ARTIFACTS[name]][1]
                written = result.message.removeprefix("Wrote ").strip()
                return (f"已生成《{title}》，写入 {written}。", result.wrote, False)
            return (result.message, [], False)
        if name == "start_teaching":
            topic = (args.get("topic") or "").strip()
            if not topic:
                return ("请说明要系统教学的主题。", [], True)
            return (self._start_teaching(topic), [], True)
        if name == "ingest_materials":
            return (self._ingest_status().message, [], False)
        if name == "show_status":
            return (self._status().message, [], False)
        if name == "lint_conflicts":
            return (self._lint().message, [], False)
        if name == "list_files":
            return (self._list_files(), [], False)
        if name == "read_file":
            return (self._read_file(str(args.get("path", ""))), [], False)
        if name == "grep_materials":
            return (self._grep_materials(str(args.get("query", ""))), [], False)
        if name == "update_memory":
            note = str(args.get("note", "")).strip()
            if not note:
                return ("没有提供要记住的内容。", [], False)
            category = str(args.get("category", "")).strip() or None
            added = self.memory.append_memory_note(note, category=category)
            if added:
                return (f"已记入长期记忆：{note}", [], False)
            return (f"这条长期记忆已存在：{note}", [], False)
        return (f"未知工具：{name}", [], False)

    def _resolve_in_workspace(self, rel_path: str) -> Path | None:
        rel = (rel_path or "").strip().strip('"').strip("'")
        if not rel:
            return None
        root = self.workspace.root.resolve()
        try:
            candidate = (root / rel).resolve() if not Path(rel).is_absolute() else Path(rel).resolve()
        except (OSError, ValueError):
            return None
        if candidate == root or root in candidate.parents:
            return candidate
        return None

    def _list_files(self) -> str:
        sources = discover_workspace_sources(self.workspace.root)
        outputs = sorted(
            (path for path in self.workspace.output_dir.rglob("*") if path.is_file()),
            key=lambda path: path.as_posix().lower(),
        )
        lines = [f"课程文件夹：{self.workspace.root}"]
        if sources:
            lines.append(f"资料（{len(sources)}）：")
            lines.extend(f"  {source.relative_path.as_posix()}  [{source.kind}]" for source in sources)
        else:
            lines.append("资料：无（把 PDF/PPT/图片/笔记放进来，再运行 /ingest）")
        if outputs:
            lines.append(f"产物（{len(outputs)}）：")
            lines.extend(f"  {path.relative_to(self.workspace.root).as_posix()}" for path in outputs)
        return "\n".join(lines)

    def _read_file(self, rel_path: str, *, max_chars: int = 8000) -> str:
        target = self._resolve_in_workspace(rel_path)
        if target is None:
            return f"路径越界或为空：{rel_path}（只能读课程文件夹内的文件）"
        if target.is_dir():
            return f"{rel_path} 是文件夹，不是文件。用 list_files 查看内容。"
        if target.exists() and target.suffix.lower() in {".md", ".txt"}:
            return self._render_file_text(target, max_chars=max_chars)
        if target.exists() and target.suffix.lower() == ".pptx":
            slides = _extract_pptx_text(target)
            if slides:
                rel = target.relative_to(self.workspace.root).as_posix()
                truncated = ""
                if len(slides) > max_chars:
                    slides = slides[:max_chars]
                    truncated = f"\n\n…（已截断，仅显示前 {max_chars} 字）"
                return f"{rel}（PPTX 文本）：\n\n{slides}{truncated}"
            # no extractable text (image-only deck) -> fall through to parsed/ingest path
        if target.exists() and target.suffix.lower() in SUPPORTED_SOURCE_EXTENSIONS:
            parsed = self._parsed_markdown_for(target)
            if parsed:
                return (
                    f"{target.relative_to(self.workspace.root).as_posix()}（MinerU 解析结果）：\n\n"
                    + self._render_file_text(parsed, max_chars=max_chars, with_header=False)
                )
            return (
                f"{target.relative_to(self.workspace.root).as_posix()} 是 {target.suffix} 文件，"
                "还没有解析。先运行 /ingest 解析后再读。"
            )
        if target.exists():
            return self._render_file_text(target, max_chars=max_chars)
        return f"找不到文件：{rel_path}"

    def _render_file_text(self, path: Path, *, max_chars: int, with_header: bool = True) -> str:
        try:
            text = path.read_text(encoding="utf-8", errors="replace")
        except OSError as exc:
            return f"读取失败：{exc}"
        truncated = ""
        if len(text) > max_chars:
            text = text[:max_chars]
            truncated = f"\n\n…（已截断，仅显示前 {max_chars} 字）"
        body = text + truncated
        if with_header:
            return f"{path.relative_to(self.workspace.root).as_posix()}：\n\n{body}"
        return body

    def _parsed_markdown_for(self, source: Path) -> Path | None:
        relative = source.relative_to(self.workspace.root)
        parsed_dir = self.workspace.cram_dir / "parsed" / "__".join(relative.with_suffix("").parts)
        if not parsed_dir.is_dir():
            return None
        markdown = sorted(parsed_dir.rglob("*.md"), key=lambda path: path.as_posix().lower())
        return markdown[0] if markdown else None

    def _grep_materials(self, query: str, *, max_hits: int = 30) -> str:
        needle = (query or "").strip()
        if not needle:
            return "请给一个要搜索的关键词。"
        lowered = needle.lower()
        hits: list[str] = []
        for path in self._text_files_for_grep():
            try:
                lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
            except OSError:
                continue
            rel = path.relative_to(self.workspace.root).as_posix()
            for number, line in enumerate(lines, start=1):
                if lowered in line.lower():
                    hits.append(f"{rel}:{number}: {line.strip()[:120]}")
                    if len(hits) >= max_hits:
                        return f"「{needle}」命中 {len(hits)}+ 处（已截断）：\n" + "\n".join(hits)
        if not hits:
            return f"没在资料里找到「{needle}」。可能需要先 /ingest 解析 PDF/PPT。"
        return f"「{needle}」命中 {len(hits)} 处：\n" + "\n".join(hits)

    def _text_files_for_grep(self) -> list[Path]:
        files = [
            source.path
            for source in discover_workspace_sources(self.workspace.root)
            if source.path.suffix.lower() in {".md", ".txt"}
        ]
        parsed = sorted(
            (self.workspace.cram_dir / "parsed").rglob("*.md"),
            key=lambda path: path.as_posix().lower(),
        )
        return files + parsed

    def _apply_switch_workspace(self, call: dict) -> tuple[str, Path | None]:
        """Switch the running router to another folder in place so the turn can continue."""
        try:
            args = json.loads(call.get("arguments") or "{}")
        except ValueError:
            args = {}
        raw = args.get("path") if isinstance(args, dict) else None
        resolved, error = self._validate_switch_path(raw)
        if resolved is None:
            return (error, None)
        self.workspace = CramWorkspace.open(resolved)
        self.memory = MemoryStore.open(self.workspace)
        return (
            f"已切换到课程文件夹：{resolved}。现在可以用 list_files 看里面有什么、用 read_file 读取文件。",
            resolved,
        )

    def _validate_switch_path(self, raw: str | None) -> tuple[Path | None, str]:
        text = (raw or "").strip().strip('"').strip("'")
        if not text:
            return None, "请说明要切换到的课程文件夹路径。"
        try:
            candidate = Path(text).expanduser()
            if not candidate.is_absolute():
                candidate = self.workspace.root / candidate
            candidate = candidate.resolve()
        except (OSError, ValueError):
            return None, f"路径无法解析：{text}"
        if not candidate.exists() or not candidate.is_dir():
            return None, f"文件夹不存在：{candidate}"
        return candidate, ""

    def _teaching_evidence(self, query: str, *, limit: int = 6) -> str:
        chunks = search_workspace_chunks(self.workspace, query, limit=limit)
        if not chunks:
            return ""
        return "\n\n".join(f"[{chunk.citation_label}]\n{chunk.text}" for chunk in chunks)

    def _start_teaching(self, topic: str) -> str:
        evidence = self._teaching_evidence(topic, limit=8)
        try:
            tree_text = self.llm.chat(
                deconstruct_messages(topic, self.workspace.course_name, evidence), stream=False
            )
        except LLMConfigurationError as exc:
            return _llm_setup_message(exc)
        except LLMRequestError as exc:
            return f"拆解知识点失败，请检查模型配置与网络。\n{exc}"
        points = parse_tree(tree_text)
        if not points:
            return f"没能拆出「{topic}」的知识点。换个更具体的主题，或先 /ingest 导入资料再试。"
        session = TeachingSession(topic=topic, points=points)
        save_session(self.workspace, session)
        return render_tree(session)

    def _teaching_turn(self, message: str):
        session = load_session(self.workspace)
        if session is None:
            yield StreamEvent("content", "教学会话已结束。直接提问即可。")
            return

        self.memory.append_session_event("user", message)
        action = classify_teaching_input(message)

        if action == "stop":
            clear_session(self.workspace)
            text = "已退出教学模式。需要时再说「教我 <主题>」，或用 /quiz 检验。"
            self.memory.append_session_event("agent", text)
            yield StreamEvent("content", text)
            return

        if action == "advance" and session.started:
            point = session.current_point()
            if point is not None:
                point.status = "taught"
            session.current += 1

        if session.finished:
            clear_session(self.workspace)
            text = (
                f"「{session.topic}」的 {len(session.points)} 个知识点都讲完了。\n"
                "输入 /quiz 出题检验，或说「教我 <新主题>」继续。"
            )
            self.memory.append_session_event("agent", text)
            yield StreamEvent("content", text)
            return

        current = session.current_point()
        evidence = self._teaching_evidence(current.title if current else session.topic, limit=5)
        messages = teach_messages(
            session, message, self.workspace.course_name, evidence, reteach=(action == "reteach")
        )
        parts: list[str] = []
        try:
            if hasattr(self.llm, "stream_chat"):
                for event in self.llm.stream_chat(messages):
                    if not isinstance(event, StreamEvent):
                        event = StreamEvent("content", str(event))
                    if event.kind == "content":
                        parts.append(event.text)
                    yield event
            else:
                text = self.llm.chat(messages, stream=False)
                parts.append(text)
                yield StreamEvent("content", text)
        except LLMConfigurationError as exc:
            text = _llm_setup_message(exc)
            self.memory.append_session_event("agent", text)
            yield StreamEvent("content", text)
            return
        except LLMRequestError as exc:
            text = "LLM 请求失败，当前会话已保留。请检查模型服务地址、模型名和网络状态。\n" + str(exc)
            self.memory.append_session_event("agent", text)
            yield StreamEvent("content", text)
            return

        session.started = True
        session.taught_since_checkin += 1
        if session.taught_since_checkin >= 3 and session.current < len(session.points) - 1:
            session.taught_since_checkin = 0
            note = "\n\n（已连讲 3 个知识点。继续就直接回应；想检验输入 /quiz；想停说「退出」。）"
            parts.append(note)
            yield StreamEvent("content", note)
        save_session(self.workspace, session)

        answer = "".join(parts).strip()
        if answer:
            self.memory.append_session_event("agent", answer)

    def _system_prompt(self) -> str:
        # Stable prefix: kept byte-identical across turns so DeepSeek's automatic prefix cache hits.
        # Per-turn retrieval is deliberately NOT here — it goes in the volatile tail (see
        # _current_user_message) so it never busts the cached prefix.
        prompt = """你是期末速成引擎，一个面向考前复习的中文学习 Agent，可以调用工具。

工作方式：
- 用户想要复习产物（速成计划/知识点整合/思维导图/题库/考前总结）时，调用对应的 generate_* 工具生成并写入文件；可用 focus 聚焦主题或题型。
- 用户想被系统讲解、带着复习某个主题（「教我X」「带我过一遍X」「系统学X」）时，调用 start_teaching 进入四步教学法。
- 导入/解析/重新索引资料用 ingest_materials；看有哪些文件、读文件、搜关键词、换课程文件夹分别用 list_files / read_file / grep_materials / switch_workspace；看状态用 show_status；查冲突用 lint_conflicts。
- 只是概念讲解、答疑、讨论时直接用中文回答，不要调用工具。
- 当你得知值得长期记住的事实（考试形式、老师强调的必考点、用户偏好、用户反复搞错或已掌握的点）时，调用 update_memory 记下来。
- 调用工具后用一两句话说明你做了什么、产物在哪，不要重复整篇内容。

回答规则：
- 用中文，结构清晰，适合考前快速复习。
- 用到资料处用方括号标注来源标签，例如 [a.pdf:mineru:1]；不要编造来源或页码。
"""
        memory_text = self.memory.load_boot_summary().strip()
        if memory_text:
            prompt += f"\n## 长期记忆\n{memory_text}\n"
        prompt += f"\n## 当前工作区\n{self._workspace_map()}\n"
        return prompt

    def _workspace_map(self, *, limit: int = 40) -> str:
        sources = discover_workspace_sources(self.workspace.root)
        outputs = sorted(
            (path for path in self.workspace.output_dir.rglob("*") if path.is_file()),
            key=lambda path: path.as_posix().lower(),
        )
        indexed = workspace_chunks_path(self.workspace).exists()
        lines = [
            f"课程：{self.workspace.course_name}（{self.workspace.root}）",
            f"资料索引：{'已建立' if indexed else '未建立（可让我 ingest_materials）'}",
        ]
        if sources:
            shown = "、".join(source.relative_path.as_posix() for source in sources[:limit])
            more = "" if len(sources) <= limit else f" 等共 {len(sources)} 个"
            lines.append(f"资料文件（{len(sources)}）：{shown}{more}")
        else:
            lines.append("资料文件：无（把 PDF/PPT/图片/笔记放进来再 ingest）")
        if outputs:
            shown = "、".join(path.relative_to(self.workspace.root).as_posix() for path in outputs[:limit])
            lines.append(f"已生成产物（{len(outputs)}）：{shown}")
        return "\n".join(lines)

    def _history_messages(self, *, exclude_last_user: str | None = None) -> list[dict]:
        events = self.memory.load_recent_session_events(limit=HISTORY_TURNS * 2)
        # run_turn/stream log the current user message before assembling context; drop that
        # trailing echo so we don't duplicate it (it's re-added as the volatile tail).
        if (
            exclude_last_user is not None
            and events
            and events[-1].get("role") == "user"
            and events[-1].get("content") == exclude_last_user
        ):
            events = events[:-1]
        messages: list[dict] = []
        for event in events:
            content = event.get("content", "")
            if not content:
                continue
            role = "assistant" if event.get("role") == "agent" else "user"
            messages.append({"role": role, "content": content})
        return messages

    def _current_user_message(self, message: str) -> dict:
        # Volatile tail: per-turn retrieval rides with the latest user message, after the stable
        # prefix and the append-only history, so it never breaks the prompt cache.
        evidence = search_workspace_chunks(self.workspace, message, limit=5)
        if evidence:
            block = "\n\n".join(f"[{chunk.citation_label}]\n{chunk.text}" for chunk in evidence)
            return {"role": "user", "content": f"{message}\n\n[本轮可参考的课程资料]\n{block}"}
        return {"role": "user", "content": message}

    def _agent_messages(self, message: str) -> list[dict]:
        messages = [{"role": "system", "content": self._system_prompt()}]
        messages.extend(self._history_messages(exclude_last_user=message))
        messages.append(self._current_user_message(message))
        return messages

    def _llm_messages(self, message: str) -> list[dict]:
        # Same cache-friendly layout as the agent path: stable system prefix + replayed history
        # + a volatile tail (current message with this turn's retrieval).
        messages = [{"role": "system", "content": self._system_prompt()}]
        messages.extend(self._history_messages(exclude_last_user=message))
        messages.append(self._current_user_message(message))
        return messages

    def _lint(self) -> CommandResult:
        conflicts = self.memory.load_conflicts()
        references = self.memory.build_reference_catalog()
        if conflicts:
            conflict_lines = "\n".join(f"- {item['title']}: {item['left']} <> {item['right']}" for item in conflicts)
        else:
            conflict_lines = "- 暂未记录冲突"
        return CommandResult(
            kind="lint",
            message=(
                "记忆健康检查\n"
                f"- 可引用条目：{len(references)}\n"
                f"- 冲突记录：{len(conflicts)}\n"
                f"{conflict_lines}"
            ),
        )


def _extract_pptx_text(path: Path) -> str:
    """Pull slide text from a .pptx quickly via python-pptx (no MinerU). Empty if unavailable."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        presentation = Presentation(str(path))
    except Exception:
        return ""
    blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    lines.append(text)
        if lines:
            blocks.append(f"# 第 {index} 页\n" + "\n".join(lines))
    return "\n\n".join(blocks)


def _default_llm_client() -> LLMClient:
    config = load_effective_llm_config()
    if config:
        return OpenAICompatibleClient(
            LLMSettings(
                provider="openai-compatible",
                base_url=config.base_url,
                model=config.model,
            ),
            api_key=config.api_key,
        )
    return OpenAICompatibleClient(_llm_settings_from_env())


def _llm_settings_from_env() -> LLMSettings:
    return LLMSettings(
        provider="openai-compatible",
        base_url=os.environ.get("CRAM_LLM_BASE_URL", "https://api.openai.com/v1"),
        model=os.environ.get("CRAM_LLM_MODEL", "gpt-4o-mini"),
        api_key_env=os.environ.get("CRAM_LLM_API_KEY_ENV", "CRAM_LLM_API_KEY"),
    )


def _format_ingest_message(
    *,
    workspace: CramWorkspace,
    total_sources: int,
    material_result: MaterialIngestResult,
    indexed_chunks: int,
    indexed_files: int,
) -> str:
    lines = [
        "资料扫描完成。",
        f"- 当前文件夹：{workspace.root}",
        f"- 找到资料：{total_sources} 个文件",
        f"- MinerU 已解析：{material_result.processed_files} 个文件",
        f"- 已建立索引：{indexed_chunks} 个片段，来自 {indexed_files} 个文本/解析结果",
    ]
    if material_result.pending_files:
        lines.append(f"- 暂未处理：{len(material_result.pending_files)} 个文件（{_preview_names(material_result.pending_files)}）")
    else:
        lines.append("- 暂未处理：0 个文件")
    if material_result.failed_files:
        lines.append(f"- 解析失败：{len(material_result.failed_files)} 个文件（{_preview_names(material_result.failed_files)}）")
    else:
        lines.append("- 解析失败：0 个文件")
    if _looks_like_code_repository(workspace.root):
        lines.extend(
            [
                "",
                "提示：你现在可能在代码仓库里运行 cram。",
                "请先进入某个学科资料文件夹，比如 D:\\期末资料\\通信原理，再运行 cram 和 /ingest。",
            ]
        )
    else:
        lines.extend(
            [
                "",
                "下一步：直接提问，或输入 /notes、/mindmap、/quiz 生成复习产物。",
            ]
        )
    return "\n".join(lines)


def _looks_like_code_repository(path: Path) -> bool:
    code_markers = {
        ".git",
        "pyproject.toml",
        "package.json",
        "requirements.txt",
        "Cargo.toml",
        "go.mod",
    }
    return any((path / marker).exists() for marker in code_markers)


def _preview_names(paths: list[str], *, limit: int = 3) -> str:
    shown = paths[:limit]
    suffix = "" if len(paths) <= limit else f" 等 {len(paths)} 个"
    return "、".join(shown) + suffix


def _llm_setup_message(error: Exception) -> str:
    return (
        "LLM 还没有配置好，当前问题已记录，但不会调用模型。\n\n"
        f"{error}\n\n"
        "在新终端里配置：\n"
        'setx CRAM_LLM_API_KEY "你的密钥"\n'
        'setx CRAM_LLM_BASE_URL "https://api.openai.com/v1"\n'
        'setx CRAM_LLM_MODEL "gpt-4o-mini"\n\n'
        "配置后重新打开终端，再进入学科资料文件夹运行 cram。"
    )


def main() -> int:
    workspace = CramWorkspace.open(Path.cwd())
    router = CommandRouter(workspace)
    print(router.handle("/status").message)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
