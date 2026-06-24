import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from app.backend.cram_app.commands import CommandRouter
from app.backend.cram_app.llm import StreamEvent
from app.backend.cram_app.teaching import (
    KnowledgePoint,
    TeachingSession,
    has_active_session,
    load_session,
    save_session,
)
from app.backend.cram_app.workspace import CramWorkspace
from app.backend.cram_app.workspace_ingest import MaterialIngestResult
from app.backend.cram_app.workspace_index import ParsedTextSource


class FakeLLM:
    def __init__(self, answer: str = "LLM 已回答"):
        self.answer = answer
        self.messages: list[dict] = []

    def chat(self, messages: list[dict], *, stream: bool = False) -> str:
        self.messages = messages
        return self.answer


class FakeStreamingLLM(FakeLLM):
    def stream_chat(self, messages: list[dict]):
        self.messages = messages
        yield "采样"
        yield "定理"


class FakeReasoningLLM(FakeLLM):
    def stream_chat(self, messages: list[dict]):
        self.messages = messages
        yield StreamEvent("reasoning", "先回忆采样定理…")
        yield StreamEvent("content", "采样")
        yield StreamEvent("content", "定理")


class FakeAgentLLM(FakeLLM):
    """Calls one tool on the first agent step, then answers on the second."""

    def __init__(self, tool_name: str | None = None, artifact_body: str = "# 题库\nQ1. 采样定理是什么？"):
        super().__init__()
        self.tool_name = tool_name
        self.artifact_body = artifact_body
        self.agent_calls = 0

    def stream_agent(self, messages: list[dict], tools: list[dict]):
        self.agent_calls += 1
        if self.agent_calls == 1 and self.tool_name:
            yield StreamEvent("tool_call", json.dumps({"id": "c1", "name": self.tool_name, "arguments": "{}"}))
        else:
            yield StreamEvent("content", "已经帮你处理好了。")

    def chat(self, messages: list[dict], *, stream: bool = False) -> str:
        # used when a generate_* tool runs the artifact generation
        self.messages = messages
        return self.artifact_body


class TuiCommandTests(unittest.TestCase):
    def test_status_reports_current_folder_sources_and_outputs(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            (workspace.root / "复习重点.pdf").write_text("pdf", encoding="utf-8")
            (workspace.output_dir / "速成计划.md").write_text("plan", encoding="utf-8")

            result = CommandRouter(workspace).handle("/status")

            self.assertEqual(result.kind, "status")
            self.assertIn("通信原理", result.message)
            self.assertIn("1 个资料文件", result.message)
            self.assertIn("1 个输出产物", result.message)

    def test_help_lists_supported_slash_commands(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))

            result = CommandRouter(workspace).handle("/help")

            self.assertEqual(result.kind, "help")
            self.assertIn("/ingest", result.message)
            self.assertIn("/config", result.message)
            self.assertIn("/model", result.message)
            self.assertIn("/lint", result.message)

    def test_model_command_routes_to_model_picker(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))

            result = CommandRouter(workspace).handle("/model")

            self.assertEqual(result.kind, "model")

    def test_generation_command_writes_llm_content(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "数字电路")
            llm = FakeLLM("# 速成计划\n1. 先掌握布尔代数化简")

            result = CommandRouter(workspace, llm=llm).handle("/plan")

            self.assertEqual(result.kind, "artifact")
            self.assertEqual(result.wrote, [workspace.output_dir / "速成计划.md"])
            written = (workspace.output_dir / "速成计划.md").read_text(encoding="utf-8")
            self.assertIn("先掌握布尔代数化简", written)  # real LLM content, not a placeholder
            self.assertNotIn("后续会接入 LLM", written)
            self.assertIn("期末速成计划", llm.messages[0]["content"])

    def test_generation_command_passes_indexed_evidence_to_llm(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            (workspace.root / "notes.md").write_text(
                "Nyquist sampling theorem prevents aliasing.",
                encoding="utf-8",
            )
            CommandRouter(workspace).handle("/ingest")
            llm = FakeLLM("# 题库\nQ1. ...")

            CommandRouter(workspace, llm=llm).handle("/quiz")

            system_content = llm.messages[0]["content"]
            self.assertIn("notes.md:text:1", system_content)
            self.assertIn("Nyquist sampling theorem", system_content)
            self.assertNotIn("当前阶段还没有接入资料检索", system_content)

    def test_generation_command_without_llm_config_does_not_write_file(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            config_path = Path(tmp) / "missing-llm.json"

            with patch.dict("os.environ", {"CRAM_LLM_CONFIG_PATH": str(config_path)}, clear=True):
                result = CommandRouter(workspace).handle("/plan")

            self.assertEqual(result.kind, "artifact")
            self.assertIn("CRAM_LLM_API_KEY", result.message)
            self.assertEqual(result.wrote, [])
            self.assertFalse((workspace.output_dir / "速成计划.md").exists())

    def test_lint_reports_existing_conflicts_and_reference_health(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "机器学习")
            router = CommandRouter(workspace)
            router.memory.record_conflict("定义冲突", left="[生成产物] A", right="[原始资料] B")

            result = router.handle("/lint")

            self.assertEqual(result.kind, "lint")
            self.assertIn("定义冲突", result.message)

    def test_free_text_is_routed_as_question(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            llm = FakeLLM("采样定理是把连续信号离散化的基础。")

            result = CommandRouter(workspace, llm=llm).handle("帮我讲一下采样定理")

            self.assertEqual(result.kind, "ask")
            self.assertEqual(result.message, "采样定理是把连续信号离散化的基础。")
            self.assertIn("期末速成", llm.messages[0]["content"])
            self.assertEqual(llm.messages[-1]["content"], "帮我讲一下采样定理")

    def test_ingest_indexes_text_sources_and_reports_pdf_as_pending(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            (workspace.root / "notes.md").write_text("Nyquist sampling theorem.", encoding="utf-8")
            (workspace.root / "slides.pdf").write_text("pdf placeholder", encoding="utf-8")

            def fake_ingest(workspace, sources):
                return MaterialIngestResult(parsed_texts=[], processed_files=0, failed_files=[], pending_files=["slides.pdf"])

            result = CommandRouter(workspace, material_ingester=fake_ingest).handle("/ingest")

            self.assertEqual(result.kind, "ingest")
            self.assertTrue((workspace.cram_dir / "index" / "chunks.jsonl").is_file())
            self.assertIn("已建立索引：1 个片段", result.message)
            self.assertIn("暂未处理：1 个文件", result.message)
            self.assertNotIn("Scanned", result.message)

    def test_ingest_indexes_mineru_markdown_from_pdf(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            (workspace.root / "slides.pdf").write_text("pdf placeholder", encoding="utf-8")
            parsed = workspace.cram_dir / "parsed" / "slides" / "auto.md"
            parsed.parent.mkdir(parents=True, exist_ok=True)
            parsed.write_text("Nyquist sampling theorem prevents aliasing.", encoding="utf-8")

            def fake_ingest(workspace, sources):
                return MaterialIngestResult(
                    parsed_texts=[
                        ParsedTextSource(path=parsed, source_file="slides.pdf", locator_prefix="mineru")
                    ],
                    processed_files=1,
                    failed_files=[],
                    pending_files=[],
                )

            result = CommandRouter(workspace, material_ingester=fake_ingest).handle("/ingest")

            self.assertIn("MinerU 已解析：1 个文件", result.message)
            self.assertIn("已建立索引：1 个片段", result.message)
            self.assertIn("slides.pdf:mineru:1", (workspace.cram_dir / "index" / "chunks.jsonl").read_text(encoding="utf-8"))

    def test_ingest_warns_when_running_inside_code_repository(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "cram-engine")
            (workspace.root / ".git").mkdir()
            (workspace.root / "README.md").write_text("# project", encoding="utf-8")

            def fake_ingest(workspace, sources):
                return MaterialIngestResult(parsed_texts=[], processed_files=0, failed_files=[], pending_files=[])

            result = CommandRouter(workspace, material_ingester=fake_ingest).handle("/ingest")

            self.assertIn("你现在可能在代码仓库里运行 cram", result.message)
            self.assertIn("请先进入某个学科资料文件夹", result.message)

    def test_free_text_question_puts_retrieval_in_volatile_tail(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            (workspace.root / "notes.md").write_text(
                "Nyquist sampling theorem prevents aliasing.",
                encoding="utf-8",
            )
            CommandRouter(workspace).handle("/ingest")
            llm = FakeLLM("use cited answer")

            CommandRouter(workspace, llm=llm).handle("What prevents aliasing?")

            # retrieval rides with the last user message (volatile tail) so the system
            # prefix stays byte-stable and cacheable
            self.assertIn("notes.md:text:1", llm.messages[-1]["content"])
            self.assertIn("Nyquist sampling theorem", llm.messages[-1]["content"])
            self.assertNotIn("notes.md:text:1", llm.messages[0]["content"])

    def test_free_text_can_stream_llm_chunks(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            llm = FakeStreamingLLM()
            router = CommandRouter(workspace, llm=llm)

            events = list(router.stream("帮我讲一下采样定理"))

            self.assertEqual([e.text for e in events if e.kind == "content"], ["采样", "定理"])
            events_log = router.memory.load_recent_session_events()
            self.assertEqual(events_log[-1]["role"], "agent")
            self.assertEqual(events_log[-1]["content"], "采样定理")

    def test_stream_surfaces_reasoning_then_content_and_saves_only_answer(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            router = CommandRouter(workspace, llm=FakeReasoningLLM())

            events = list(router.stream("帮我讲一下采样定理"))

            self.assertEqual(
                [(e.kind, e.text) for e in events],
                [("reasoning", "先回忆采样定理…"), ("content", "采样"), ("content", "定理")],
            )
            saved = router.memory.load_recent_session_events()[-1]
            self.assertEqual(saved["role"], "agent")
            self.assertEqual(saved["content"], "采样定理")  # reasoning is not persisted

    def test_free_text_reports_missing_llm_api_key_without_crashing(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            config_path = Path(tmp) / "missing-llm.json"

            with patch.dict("os.environ", {"CRAM_LLM_CONFIG_PATH": str(config_path)}, clear=True):
                result = CommandRouter(workspace).handle("帮我讲一下采样定理")

            self.assertEqual(result.kind, "ask")
            self.assertIn("CRAM_LLM_API_KEY", result.message)
            self.assertIn("setx", result.message)

    def test_run_turn_invokes_tool_then_answers(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            llm = FakeAgentLLM(tool_name="generate_quiz")
            router = CommandRouter(workspace, llm=llm)

            events = list(router.run_turn("帮我出一套题"))

            kinds = [event.kind for event in events]
            self.assertIn("tool", kinds)
            self.assertIn("wrote", kinds)
            self.assertIn("content", kinds)
            # the quiz artifact was actually generated and written
            self.assertTrue((workspace.output_dir / "题库.md").is_file())
            written = (workspace.output_dir / "题库.md").read_text(encoding="utf-8")
            self.assertIn("采样定理是什么", written)
            # the final answer is persisted to memory
            self.assertEqual(router.memory.load_recent_session_events()[-1]["role"], "agent")

    def test_run_turn_without_tool_just_streams_answer(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            router = CommandRouter(workspace, llm=FakeAgentLLM(tool_name=None))

            events = list(router.run_turn("讲一下采样定理"))

            self.assertEqual([e.text for e in events if e.kind == "content"], ["已经帮你处理好了。"])
            self.assertEqual([e.kind for e in events if e.kind == "tool"], [])

    def test_run_turn_falls_back_to_plain_stream_without_tool_support(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp))
            router = CommandRouter(workspace, llm=FakeStreamingLLM())

            events = list(router.run_turn("讲一下采样定理"))

            self.assertEqual([e.text for e in events if e.kind == "content"], ["采样", "定理"])

    def test_run_turn_passes_focus_to_generated_artifact(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")

            class FocusAgentLLM(FakeAgentLLM):
                def stream_agent(self, messages, tools):
                    self.agent_calls += 1
                    if self.agent_calls == 1:
                        yield StreamEvent(
                            "tool_call",
                            json.dumps({"id": "c1", "name": "generate_quiz", "arguments": '{"focus": "傅里叶变换"}'}),
                        )
                    else:
                        yield StreamEvent("content", "已按傅里叶变换出题。")

            llm = FocusAgentLLM()
            router = CommandRouter(workspace, llm=llm)

            list(router.run_turn("只出傅里叶变换的题"))

            # the artifact generation prompt received the focus
            self.assertIn("傅里叶变换", llm.messages[0]["content"])


class FakeTeachLLM(FakeLLM):
    """stream_agent calls start_teaching; chat returns the deconstruct tree; stream_chat teaches."""

    def __init__(self, tree_json='{"points": [{"title": "采样定理", "hook": "none"}, {"title": "频谱混叠", "hook": "contrast"}]}'):
        super().__init__()
        self.tree_json = tree_json

    def stream_agent(self, messages, tools):
        yield StreamEvent("tool_call", json.dumps({"id": "c1", "name": "start_teaching", "arguments": '{"topic": "采样定理"}'}))

    def chat(self, messages, *, stream=False):
        self.messages = messages
        return self.tree_json

    def stream_chat(self, messages):
        self.messages = messages
        yield StreamEvent("content", "先想象你在录音…这就是采样定理。")


class TeachingTurnTests(unittest.TestCase):
    def test_run_turn_start_teaching_creates_session_and_shows_tree(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeTeachLLM())

            events = list(router.run_turn("教我采样定理"))

            kinds = [e.kind for e in events]
            content = "".join(e.text for e in events if e.kind == "content")
            self.assertIn("tool", kinds)  # start_teaching surfaced as a tool action
            self.assertIn("采样定理", content)  # the tree topic is shown directly
            self.assertIn("频谱混叠", content)  # a deconstructed point
            self.assertTrue(has_active_session(workspace))

    def test_active_session_teaches_then_advances(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            save_session(
                workspace,
                TeachingSession(topic="采样定理", points=[KnowledgePoint("采样"), KnowledgePoint("混叠")]),
            )
            router = CommandRouter(workspace, llm=FakeTeachLLM())

            list(router.run_turn("开始"))  # teach point 0
            after_first = load_session(workspace)
            self.assertTrue(after_first.started)
            self.assertEqual(after_first.current, 0)

            list(router.run_turn("懂了，继续"))  # advance to point 1
            after_second = load_session(workspace)
            self.assertEqual(after_second.current, 1)
            self.assertEqual(after_second.points[0].status, "taught")

    def test_active_session_stop_clears_session(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            save_session(workspace, TeachingSession(topic="采样定理", points=[KnowledgePoint("采样")]))
            router = CommandRouter(workspace, llm=FakeTeachLLM())

            events = list(router.run_turn("退出"))

            self.assertFalse(has_active_session(workspace))
            self.assertIn("退出", "".join(e.text for e in events if e.kind == "content"))

    def test_teaching_finish_hands_off_to_quiz(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            save_session(
                workspace,
                TeachingSession(topic="采样定理", points=[KnowledgePoint("采样")], started=True, current=0),
            )
            router = CommandRouter(workspace, llm=FakeTeachLLM())

            events = list(router.run_turn("懂了"))  # advance past the last point -> finished

            content = "".join(e.text for e in events if e.kind == "content")
            self.assertIn("/quiz", content)
            self.assertFalse(has_active_session(workspace))


class FakeSwitchLLM(FakeLLM):
    def __init__(self, target: str):
        super().__init__()
        self.target = target
        self.calls = 0

    def stream_agent(self, messages, tools):
        self.calls += 1
        if self.calls == 1:
            yield StreamEvent(
                "tool_call",
                json.dumps({"id": "c1", "name": "switch_workspace", "arguments": json.dumps({"path": self.target})}),
            )
        else:
            yield StreamEvent("content", "好的，已经切换。")


class FakeSwitchThenListLLM(FakeLLM):
    """Switch folders, then keep going (list files), then answer — a multi-step turn."""

    def __init__(self, target: str):
        super().__init__()
        self.target = target
        self.calls = 0

    def stream_agent(self, messages, tools):
        self.calls += 1
        if self.calls == 1:
            yield StreamEvent(
                "tool_call",
                json.dumps({"id": "c1", "name": "switch_workspace", "arguments": json.dumps({"path": self.target})}),
            )
        elif self.calls == 2:
            yield StreamEvent("tool_call", json.dumps({"id": "c2", "name": "list_files", "arguments": "{}"}))
        else:
            yield StreamEvent("content", "已切换并查看了文件。")


class FileToolTests(unittest.TestCase):
    def _router(self, root: Path) -> CommandRouter:
        return CommandRouter(CramWorkspace.open(root), llm=FakeLLM())

    def test_list_files_shows_sources_and_outputs(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            (workspace.root / "notes.md").write_text("hi", encoding="utf-8")
            (workspace.output_dir / "速成计划.md").write_text("plan", encoding="utf-8")
            router = CommandRouter(workspace, llm=FakeLLM())

            text, wrote, direct = router._execute_tool({"name": "list_files", "arguments": "{}"})

            self.assertIn("notes.md", text)
            self.assertIn("cram-output/速成计划.md", text)
            self.assertEqual(wrote, [])
            self.assertFalse(direct)

    def test_read_file_reads_text_within_workspace(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            (workspace.root / "notes.md").write_text("奈奎斯特采样定理的内容", encoding="utf-8")
            router = CommandRouter(workspace, llm=FakeLLM())

            text, _, _ = router._execute_tool({"name": "read_file", "arguments": '{"path": "notes.md"}'})

            self.assertIn("奈奎斯特采样定理的内容", text)

    def test_read_file_blocks_path_traversal(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            (Path(tmp) / "secret.txt").write_text("机密", encoding="utf-8")
            router = CommandRouter(workspace, llm=FakeLLM())

            text, _, _ = router._execute_tool({"name": "read_file", "arguments": '{"path": "../secret.txt"}'})

            self.assertIn("越界", text)
            self.assertNotIn("机密", text)

    def test_read_unparsed_pdf_suggests_ingest(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            (workspace.root / "slides.pdf").write_text("pdf bytes", encoding="utf-8")
            router = CommandRouter(workspace, llm=FakeLLM())

            text, _, _ = router._execute_tool({"name": "read_file", "arguments": '{"path": "slides.pdf"}'})

            self.assertIn("/ingest", text)

    def test_grep_materials_returns_file_and_line(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            (workspace.root / "a.md").write_text("第一行\n采样定理很重要\n结尾", encoding="utf-8")
            router = CommandRouter(workspace, llm=FakeLLM())

            text, _, _ = router._execute_tool({"name": "grep_materials", "arguments": '{"query": "采样定理"}'})

            self.assertIn("a.md:2", text)
            self.assertIn("采样定理很重要", text)

    def test_run_turn_switch_workspace_emits_switch_event(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            other = Path(tmp) / "数字电路"
            other.mkdir()
            router = CommandRouter(workspace, llm=FakeSwitchLLM(str(other)))

            events = list(router.run_turn("切换到数字电路那门课"))

            switch_events = [e for e in events if e.kind == "switch"]
            self.assertEqual(len(switch_events), 1)
            self.assertEqual(Path(switch_events[0].text).name, "数字电路")

    def test_run_turn_continues_after_switch_with_more_tools(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            other = Path(tmp) / "数字电路"
            other.mkdir()
            (other / "讲义.md").write_text("数字电路讲义", encoding="utf-8")
            router = CommandRouter(workspace, llm=FakeSwitchThenListLLM(str(other)))

            events = list(router.run_turn("切到数字电路看看有什么文件"))

            # switched in place, and the turn did NOT stop after the switch
            self.assertEqual(router.workspace.course_name, "数字电路")
            self.assertEqual(len([e for e in events if e.kind == "switch"]), 1)
            self.assertGreaterEqual(len([e for e in events if e.kind == "tool"]), 2)
            content = "".join(e.text for e in events if e.kind == "content")
            self.assertIn("已切换并查看了文件", content)

    def test_read_pptx_extracts_slide_text_without_ingest(self):
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # blank
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            title_box.text_frame.text = "采样定理"
            body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
            body_box.text_frame.text = "离散采样无失真还原连续信号"
            presentation.save(str(workspace.root / "lecture.pptx"))
            router = CommandRouter(workspace, llm=FakeLLM())

            text, _, _ = router._execute_tool({"name": "read_file", "arguments": '{"path": "lecture.pptx"}'})

            self.assertIn("采样定理", text)
            self.assertIn("离散采样无失真还原连续信号", text)
            self.assertNotIn("/ingest", text)  # fast python-pptx path, no MinerU needed

    def test_validate_switch_path_rejects_missing_dir(self):
        with tempfile.TemporaryDirectory() as tmp:
            router = CommandRouter(CramWorkspace.open(Path(tmp) / "通信原理"), llm=FakeLLM())

            resolved, error = router._validate_switch_path(str(Path(tmp) / "nope"))

            self.assertIsNone(resolved)
            self.assertIn("不存在", error)


class FakeLoopingLLM(FakeLLM):
    """Always calls a tool and never answers, to exercise the agent step cap."""

    def stream_agent(self, messages, tools):
        yield StreamEvent("tool_call", json.dumps({"id": "c1", "name": "show_status", "arguments": "{}"}))

    def stream_chat(self, messages):
        yield StreamEvent("content", "已尽力检索，先这样。")


class AgentLoopGuardTests(unittest.TestCase):
    def test_run_turn_forces_final_answer_when_tool_loop_caps_out(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLoopingLLM())

            events = list(router.run_turn("一直调工具看看"))

            contents = [e.text for e in events if e.kind == "content"]
            self.assertIn("已尽力检索，先这样。", contents)  # forced closing answer, not silence
            self.assertEqual(router.memory.load_recent_session_events()[-1]["role"], "agent")


class FakeUpdateMemoryLLM(FakeLLM):
    """Records a durable fact via update_memory, then answers."""

    def __init__(self, note: str, category: str = "考点"):
        super().__init__()
        self.note = note
        self.category = category
        self.calls = 0

    def stream_agent(self, messages, tools):
        self.calls += 1
        if self.calls == 1:
            yield StreamEvent(
                "tool_call",
                json.dumps(
                    {
                        "id": "m1",
                        "name": "update_memory",
                        "arguments": json.dumps({"note": self.note, "category": self.category}),
                    }
                ),
            )
        else:
            yield StreamEvent("content", "记住了。")


class LongTermMemoryTests(unittest.TestCase):
    def test_update_memory_tool_persists_and_feeds_next_turn(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeUpdateMemoryLLM("第3章必考"))

            list(router.run_turn("老师说第3章必考"))

            # persisted to long-term memory
            self.assertIn("第3章必考", router.memory.load_boot_summary())
            # and now rides in the stable prefix on the next turn
            self.assertIn("第3章必考", router._system_prompt())
            self.assertIn("长期记忆", router._system_prompt())
    def test_system_prompt_includes_workspace_map_and_excludes_retrieval(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            (workspace.root / "notes.md").write_text(
                "Nyquist sampling theorem prevents aliasing.", encoding="utf-8"
            )
            router = CommandRouter(workspace, llm=FakeLLM())
            router.handle("/ingest")

            system = router._system_prompt()

            self.assertIn("通信原理", system)  # workspace map present
            self.assertIn("notes.md", system)  # file listed in the map
            self.assertNotIn("Nyquist sampling theorem", system)  # raw chunk text not in stable prefix

    def test_agent_messages_replays_recent_history(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLLM())
            router.memory.append_session_event("user", "采样定理是什么")
            router.memory.append_session_event("agent", "它说的是……")

            messages = router._agent_messages("再讲讲混叠")

            self.assertEqual([m["role"] for m in messages], ["system", "user", "assistant", "user"])
            self.assertEqual(messages[1]["content"], "采样定理是什么")
            self.assertEqual(messages[2]["content"], "它说的是……")
            self.assertTrue(messages[-1]["content"].startswith("再讲讲混叠"))

    def test_system_prompt_stable_across_conversation_turns(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLLM())

            first = router._system_prompt()
            router.memory.append_session_event("user", "hi")
            router.memory.append_session_event("agent", "hello")
            second = router._system_prompt()

            self.assertEqual(first, second)  # cacheable prefix must not change per turn

    def test_history_excludes_current_user_echo(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLLM())
            router.memory.append_session_event("user", "当前这句话")  # simulate run_turn pre-logging

            messages = router._agent_messages("当前这句话")

            # the just-logged current message appears once (as the tail), not duplicated in history
            user_contents = [m["content"] for m in messages if m["role"] == "user"]
            self.assertEqual(user_contents.count("当前这句话"), 1)


class CompactionTests(unittest.TestCase):
    def test_compaction_folds_aged_out_turns_into_rolling_summary(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            llm = FakeLLM("摘要：覆盖了前面几轮复习。")
            router = CommandRouter(workspace, llm=llm)
            for i in range(15):  # 30 events; 30-20 keep = 10 aged out >= COMPACT_MIN_FOLD
                router.memory.append_session_event("user", f"问题{i}")
                router.memory.append_session_event("agent", f"回答{i}")

            router._maybe_compact()

            self.assertIn("摘要", router.memory.load_rolling_summary())
            self.assertEqual(router.memory.load_summarized_through(), 10)
            messages = router._assemble_messages("新问题")
            self.assertTrue(
                any(m["role"] == "system" and "之前对话摘要" in m["content"] for m in messages)
            )

    def test_no_compaction_for_short_conversation(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLLM("不该被调用"))
            for i in range(4):
                router.memory.append_session_event("user", f"问题{i}")
                router.memory.append_session_event("agent", f"回答{i}")

            router._maybe_compact()

            self.assertEqual(router.memory.load_rolling_summary(), "")
            self.assertEqual(router.memory.load_summarized_through(), 0)


class RenderCommandTests(unittest.TestCase):
    def _decode(self, html: str) -> str:
        import base64
        import re

        token = re.search(r'atob\("([^"]+)"\)', html).group(1)
        return base64.b64decode(token).decode("utf-8")

    def test_render_builds_html_for_last_answer(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLLM())
            router.memory.append_session_event("user", "讲采样定理")
            router.memory.append_session_event("agent", r"采样定理：$f_s \geq 2f_{max}$。")

            result = router.handle("/render")

            self.assertEqual(result.kind, "render")
            html_path = workspace.cram_dir / "render" / "latest.html"
            self.assertEqual(result.wrote, [html_path])
            html = html_path.read_text(encoding="utf-8")
            self.assertIn("katex", html)  # real LaTeX rendering
            self.assertIn("texmath", html)
            decoded = self._decode(html)
            self.assertIn("f_s", decoded)  # math kept as LaTeX for the browser to render
            self.assertIn("采样定理", decoded)

    def test_render_all_includes_both_roles(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLLM())
            router.memory.append_session_event("user", "问题甲")
            router.memory.append_session_event("agent", "回答乙")

            result = router.handle("/render all")

            decoded = self._decode((workspace.cram_dir / "render" / "latest.html").read_text(encoding="utf-8"))
            self.assertIn("问题甲", decoded)
            self.assertIn("回答乙", decoded)

    def test_render_without_answer_reports_message(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            router = CommandRouter(workspace, llm=FakeLLM())

            result = router.handle("/render")

            self.assertEqual(result.kind, "render")
            self.assertIn("还没有可渲染", result.message)


class FakeMindmapLLM(FakeLLM):
    """First reply is an invalid flat line; the second is a valid outline — exercises repair."""

    def __init__(self):
        super().__init__()
        self.calls = 0

    def chat(self, messages, *, stream=False):
        self.messages = messages
        self.calls += 1
        if self.calls == 1:
            return "就这一行没有结构"
        return "# 采样定理\n- 核心结论\n  - 采样率≥2倍\n- 混叠\n  - 抗混叠滤波\n"


class MindmapCommandTests(unittest.TestCase):
    def test_mindmap_validates_repairs_and_writes_three_formats(self):
        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")
            llm = FakeMindmapLLM()

            result = CommandRouter(workspace, llm=llm).handle("/mindmap")

            self.assertEqual(result.kind, "mindmap")
            self.assertGreaterEqual(llm.calls, 2)  # invalid first attempt was repaired
            out = workspace.output_dir
            self.assertTrue((out / "思维导图.html").is_file())
            self.assertTrue((out / "思维导图.opml").is_file())
            self.assertTrue((out / "思维导图.md").is_file())
            html = (out / "思维导图.html").read_text(encoding="utf-8")
            self.assertIn("markmap-autoloader", html)
            self.assertIn("采样定理", html)
            self.assertEqual({p.suffix for p in result.wrote}, {".html", ".opml", ".md"})

    def test_mindmap_gives_up_after_repeated_invalid_output(self):
        class AlwaysBad(FakeLLM):
            def chat(self, messages, *, stream=False):
                return "一行"

        with tempfile.TemporaryDirectory() as tmp:
            workspace = CramWorkspace.open(Path(tmp) / "通信原理")

            result = CommandRouter(workspace, llm=AlwaysBad()).handle("/mindmap")

            self.assertEqual(result.kind, "artifact")  # reported failure, not a mindmap
            self.assertIn("校验未通过", result.message)
            self.assertFalse((workspace.output_dir / "思维导图.html").exists())


if __name__ == "__main__":
    unittest.main()
